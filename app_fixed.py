
from fastapi import FastAPI, Request, Form, UploadFile, File
from fastapi.responses import HTMLResponse, JSONResponse, RedirectResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from pathlib import Path
from collections import defaultdict
import json, csv, io, datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

BASE_DIR = Path(__file__).resolve().parent
DATA_FILE = BASE_DIR / "data" / "use_cases.json"
SETTINGS_FILE = BASE_DIR / "data" / "settings.json"

app = FastAPI(title="ASOC AI ME Dashboard", version="5.0.0")
app.mount("/static", StaticFiles(directory=str(BASE_DIR / "static")), name="static")
templates = Jinja2Templates(directory=str(BASE_DIR / "templates"))

FIELDS = [
    "id","name","domain","description","strategic_outcome","benefit_type","kpi",
    "baseline_kpi","current_kpi","target_kpi","benefit_confidence",
    "market","technical_owner","sponsor","status","tmf_level","impact","frequency",
    "annual_value","implementation_cost","hours_saved","mttr_improvement","reuse_score"
]

REUSE_LEGEND = [
    {"score": 1, "level": "Team Reuse", "meaning": "Works for one team or local process only."},
    {"score": 2, "level": "Function Reuse", "meaning": "Reusable within one ASOC function."},
    {"score": 3, "level": "Domain Reuse", "meaning": "Reusable across an entire FM, PM, CM or OSS domain."},
    {"score": 4, "level": "Cross-Domain Reuse", "meaning": "Reusable across multiple ASOC domains."},
    {"score": 5, "level": "Cross-Market Reuse", "meaning": "Build once, deploy across ZA, TZ, MZ, Lesotho and potentially Vodafone Group."},
]

def load_json(path, default):
    if not path.exists():
        return default
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

def save_json(path, data):
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)

def load_use_cases():
    return load_json(DATA_FILE, [])

def save_use_cases(data):
    save_json(DATA_FILE, data)

def load_settings():
    settings = load_json(SETTINGS_FILE, {})
    settings.setdefault("active_year", "FY27")
    settings.setdefault("targets", {})
    for year in ["FY27","FY28","FY29"]:
        settings["targets"].setdefault(year, {})
    return settings

def as_number(v, default=0):
    try:
        if v in [None, ""]:
            return default
        return float(v) if "." in str(v) else int(v)
    except Exception:
        return default

def safe_div(a, b):
    a = as_number(a); b = as_number(b)
    return 0 if b == 0 else a / b

def summary(use_cases, settings=None):
    settings = settings or load_settings()
    total_value = sum(as_number(x.get("annual_value", 0)) for x in use_cases)
    implementation_cost = sum(as_number(x.get("implementation_cost", 0)) for x in use_cases)
    total_hours = sum(as_number(x.get("hours_saved", 0)) for x in use_cases)
    production = sum(1 for x in use_cases if x.get("status") == "Production")
    sandbox = sum(1 for x in use_cases if x.get("status") == "Sandbox")
    prototype = sum(1 for x in use_cases if x.get("status") == "Prototype")
    in_progress = sum(1 for x in use_cases if x.get("status") == "In Progress")
    reusable_count = sum(1 for x in use_cases if as_number(x.get("reuse_score", 0)) >= 4)
    reusable_percent = round(safe_div(reusable_count, len(use_cases)) * 100, 1)
    roi = round(safe_div(total_value - implementation_cost, implementation_cost) * 100, 1) if implementation_cost else 0
    return {
        "active_use_cases": len(use_cases),
        "production": production,
        "sandbox": sandbox,
        "prototype": prototype,
        "in_progress": in_progress,
        "annual_value": total_value,
        "implementation_cost": implementation_cost,
        "hours_saved": total_hours,
        "avg_mttr": round(sum(as_number(x.get("mttr_improvement", 0)) for x in use_cases) / max(len(use_cases), 1), 1),
        "avg_reuse": round(sum(as_number(x.get("reuse_score", 0)) for x in use_cases) / max(len(use_cases), 1), 1),
        "reusable_count": reusable_count,
        "reusable_percent": reusable_percent,
        "portfolio_roi": roi,
        "tmf_current": as_number(settings.get("tmf_current", 1.8)),
        "tmf_target": as_number(settings.get("tmf_target", 3.0)),
    }

def active_targets(settings):
    active_year = settings.get("active_year", "FY27")
    targets = settings.get("targets", {}).get(active_year, {})
    return active_year, targets

def pct(current, target):
    return round(min(safe_div(current, target) * 100, 100), 1)

def gauge_values(s, settings):
    year, targets = active_targets(settings)
    return {
        "year": year,
        "annual_benefits": {
            "label": "Annual benefits realised",
            "current": s["annual_value"],
            "target": as_number(targets.get("annual_benefits", 1)),
            "unit": "R",
            "pct": pct(s["annual_value"], targets.get("annual_benefits", 1))
        },
        "hours_saved": {
            "label": "Hours returned",
            "current": s["hours_saved"],
            "target": as_number(targets.get("hours_saved", 1)),
            "unit": "hrs",
            "pct": pct(s["hours_saved"], targets.get("hours_saved", 1))
        },
        "production_use_cases": {
            "label": "Production adoption",
            "current": s["production"],
            "target": as_number(targets.get("production_use_cases", 1)),
            "unit": "use cases",
            "pct": pct(s["production"], targets.get("production_use_cases", 1))
        },
        "tmf_maturity": {
            "label": "TMF maturity",
            "current": s["tmf_current"],
            "target": as_number(targets.get("tmf_maturity", settings.get("tmf_target", 3.0))),
            "unit": "level",
            "pct": pct(s["tmf_current"], targets.get("tmf_maturity", settings.get("tmf_target", 3.0)))
        },
        "reusable_percent": {
            "label": "Cross-domain / market reusable",
            "current": s["reusable_percent"],
            "target": as_number(targets.get("reusable_use_cases_percent", 1)),
            "unit": "%",
            "pct": pct(s["reusable_percent"], targets.get("reusable_use_cases_percent", 1))
        },
        "roi": {
            "label": "Portfolio ROI",
            "current": s["portfolio_roi"],
            "target": as_number(targets.get("portfolio_roi_percent", 1)),
            "unit": "%",
            "pct": pct(s["portfolio_roi"], targets.get("portfolio_roi_percent", 1))
        }
    }

def group_by(use_cases, key, measure):
    buckets = defaultdict(float)
    for x in use_cases:
        buckets[x.get(key, "Unknown") or "Unknown"] += as_number(x.get(measure, 0))
    return [{"label": k, "value": v} for k, v in sorted(buckets.items(), key=lambda i: i[1], reverse=True)]

@app.get("/", response_class=HTMLResponse)
async def dashboard(request: Request):
    cases = load_use_cases()
    settings = load_settings()
    s = summary(cases, settings)
    return templates.TemplateResponse("index.html", {
        "request": request,
        "summary": s,
        "gauges": gauge_values(s, settings),
        "use_cases": cases,
        "settings": settings,
        "reuse_legend": REUSE_LEGEND
    })

@app.get("/admin", response_class=HTMLResponse)
async def admin(request: Request, edit: str = ""):
    cases = load_use_cases()
    settings = load_settings()
    record = next((x for x in cases if x.get("id") == edit), None)
    return templates.TemplateResponse("admin.html", {
        "request": request,
        "use_cases": cases,
        "settings": settings,
        "record": record,
        "fields": FIELDS,
        "reuse_legend": REUSE_LEGEND
    })

@app.post("/admin/save-use-case")
async def save_use_case(request: Request):
    form = await request.form()
    cases = load_use_cases()
    item = {}
    for field in FIELDS:
        value = form.get(field, "")
        if field in ["annual_value","implementation_cost","hours_saved","mttr_improvement","reuse_score"]:
            item[field] = as_number(value)
        else:
            item[field] = str(value).strip()
    if not item.get("id"):
        return RedirectResponse("/admin", status_code=303)
    found = False
    for i, c in enumerate(cases):
        if c.get("id") == item["id"]:
            cases[i] = item
            found = True
            break
    if not found:
        cases.append(item)
    save_use_cases(cases)
    return RedirectResponse("/admin", status_code=303)

@app.post("/admin/delete-use-case")
async def delete_use_case(id: str = Form(...)):
    cases = [x for x in load_use_cases() if x.get("id") != id]
    save_use_cases(cases)
    return RedirectResponse("/admin", status_code=303)

@app.post("/admin/save-settings")
async def save_settings(request: Request):
    form = await request.form()
    settings = load_settings()
    settings["dashboard_title"] = str(form.get("dashboard_title", settings.get("dashboard_title","ME Executive Dashboard")))
    settings["portfolio_period"] = str(form.get("portfolio_period", settings.get("portfolio_period","FY27–FY29")))
    settings["active_year"] = str(form.get("active_year", settings.get("active_year","FY27")))
    settings["tmf_current"] = as_number(form.get("tmf_current", settings.get("tmf_current", 1.8)))
    settings["tmf_target"] = as_number(form.get("tmf_target", settings.get("tmf_target", 3.0)))

    for year in ["FY27","FY28","FY29"]:
        settings.setdefault("targets", {}).setdefault(year, {})
        for metric in ["annual_benefits","hours_saved","production_use_cases","tmf_maturity","reusable_use_cases_percent","portfolio_roi_percent"]:
            settings["targets"][year][metric] = as_number(form.get(f"{year}_{metric}", settings["targets"][year].get(metric, 0)))

    decisions = str(form.get("executive_decisions", ""))
    settings["executive_decisions"] = [x.strip() for x in decisions.splitlines() if x.strip()]

    parsed_bets = []
    big_bets = str(form.get("big_bets", ""))
    for line in big_bets.splitlines():
        if "|" in line:
            name, desc = line.split("|", 1)
            parsed_bets.append({"name": name.strip(), "description": desc.strip()})
        elif line.strip():
            parsed_bets.append({"name": line.strip(), "description": ""})
    settings["big_bets"] = parsed_bets
    save_json(SETTINGS_FILE, settings)
    return RedirectResponse("/admin", status_code=303)

@app.get("/export/json")
async def export_json():
    data = {"use_cases": load_use_cases(), "settings": load_settings()}
    payload = json.dumps(data, indent=2).encode("utf-8")
    return StreamingResponse(io.BytesIO(payload), media_type="application/json", headers={"Content-Disposition":"attachment; filename=asoc_ai_dashboard_data.json"})

@app.get("/export/csv")
async def export_csv():
    output = io.StringIO()
    writer = csv.DictWriter(output, fieldnames=FIELDS)
    writer.writeheader()
    for row in load_use_cases():
        writer.writerow({k: row.get(k, "") for k in FIELDS})
    return StreamingResponse(io.BytesIO(output.getvalue().encode("utf-8")), media_type="text/csv", headers={"Content-Disposition":"attachment; filename=asoc_ai_use_cases.csv"})

@app.post("/import/csv")
async def import_csv(file: UploadFile = File(...)):
    content = (await file.read()).decode("utf-8-sig")
    reader = csv.DictReader(io.StringIO(content))
    existing_cases = load_use_cases()
    existing_by_id = {str(item.get("id", "")).strip(): item for item in existing_cases if str(item.get("id", "")).strip()}
    imported_by_id = {}
    for row in reader:
        item = {k: row.get(k, "") for k in FIELDS}
        item_id = str(item.get("id", "")).strip()
        if not item_id:
            continue
        item["id"] = item_id
        for n in ["annual_value","implementation_cost","hours_saved","mttr_improvement","reuse_score"]:
            item[n] = as_number(item.get(n, 0))
        imported_by_id[item_id] = item
    existing_by_id.update(imported_by_id)
    existing_order = [str(item.get("id", "")).strip() for item in existing_cases if str(item.get("id", "")).strip()]
    new_ids = [item_id for item_id in imported_by_id.keys() if item_id not in existing_order]
    final_ids = existing_order + new_ids
    merged_cases = [existing_by_id[item_id] for item_id in final_ids if item_id in existing_by_id]
    save_use_cases(merged_cases)
    return RedirectResponse("/admin", status_code=303)

@app.get("/api/use-cases")
async def api_use_cases():
    return JSONResponse(load_use_cases())

@app.get("/api/summary")
async def api_summary():
    cases = load_use_cases()
    settings = load_settings()
    s = summary(cases, settings)
    return JSONResponse({
        "summary": s,
        "gauges": gauge_values(s, settings),
        "value_by_domain": group_by(cases, "domain", "annual_value"),
        "hours_by_domain": group_by(cases, "domain", "hours_saved"),
        "value_by_sponsor": group_by(cases, "sponsor", "annual_value"),
        "status_mix": group_by(cases, "status", "annual_value"),
        "reuse_mix": group_by(cases, "reuse_score", "annual_value")
    })

def add_title(slide, title, subtitle=None):
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.5))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.42), Inches(0.78), Inches(12.2), Inches(0.3))
        p2 = st.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(220,220,220)

def set_bg(slide, color=(17,18,23)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)

def add_kpi(slide, x, y, w, h, label, value, note=""):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(34,37,48)
    shape.line.color.rgb = RGBColor(70,70,80)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(180,180,185)
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255,255,255)
    if note:
        p3 = tf.add_paragraph()
        p3.text = note
        p3.font.size = Pt(8)
        p3.font.color.rgb = RGBColor(170,170,175)

def make_bar_chart(data, title, ylabel, path):
    labels = [str(x["label"]) for x in data]
    vals = [x["value"] for x in data]
    fig, ax = plt.subplots(figsize=(7, 3.8), dpi=160)
    fig.patch.set_facecolor("#1a1c23")
    ax.set_facecolor("#1a1c23")
    ax.bar(labels, vals)
    ax.set_title(title, color="white", fontsize=12, pad=10)
    ax.set_ylabel(ylabel, color="white", fontsize=9)
    ax.tick_params(colors="white", labelsize=8)
    for spine in ax.spines.values():
        spine.set_color("#444")
    ax.grid(axis="y", color="white", alpha=.08)
    plt.xticks(rotation=20, ha="right")
    plt.tight_layout()
    fig.savefig(path, facecolor=fig.get_facecolor(), bbox_inches="tight")
    plt.close(fig)

def make_gauge(current, target, title, unit, path):
    current = float(current or 0)
    target = float(target or 1)
    pct_raw = (current / target) * 100 if target else 0
    pct_value = max(0, min(pct_raw / 100, 1))

    fig, ax = plt.subplots(figsize=(4.2, 2.8), dpi=180, subplot_kw={"projection": "polar"})
    fig.patch.set_facecolor("#1a1c23")
    ax.set_facecolor("#1a1c23")

    theta_bg = np.linspace(np.pi, 0, 100)
    ax.plot(theta_bg, [1]*len(theta_bg), linewidth=18, solid_capstyle="round", color="#e60000")

    theta_val = np.linspace(np.pi, np.pi - pct_value*np.pi, 100)
    ax.plot(theta_val, [1]*len(theta_val), linewidth=18, solid_capstyle="round", color="#22c55e")

    needle_theta = np.pi - pct_value*np.pi
    ax.plot([needle_theta, needle_theta], [0, .95], linewidth=2.5, color="white")
    ax.scatter([needle_theta], [0], s=30, color="white")

    ax.set_ylim(0, 1.25)
    ax.set_axis_off()

    current_label = f"{current:,.1f}" if current % 1 else f"{current:,.0f}"
    target_label = f"{target:,.1f}" if target % 1 else f"{target:,.0f}"
    pct_label = f"{pct_raw:,.0f}%"

    ax.text(.5, .42, pct_label, transform=ax.transAxes, ha="center", va="center",
            color="white", fontsize=18, fontweight="bold")
    ax.text(.5, .27, title, transform=ax.transAxes, ha="center", va="center",
            color="#d4d4d8", fontsize=9, fontweight="bold")
    ax.text(.5, .13, f"Current: {current_label}{unit}  |  Target: {target_label}{unit}",
            transform=ax.transAxes, ha="center", va="center", color="#a1a1aa", fontsize=7.5)

    plt.tight_layout(pad=.3)
    fig.savefig(path, facecolor=fig.get_facecolor(), bbox_inches="tight")
    plt.close(fig)

@app.get("/export/pptx")
async def export_pptx():
    cases = load_use_cases()
    settings = load_settings()
    s = summary(cases, settings)
    gauges = gauge_values(s, settings)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "ASOC AI ME Executive Dashboard", f"{settings.get('active_year','FY27')} · Exported {datetime.date.today().isoformat()}")
    kpis = [
        ("Active AI use cases", str(s["active_use_cases"]), "FM · PM · CM · OSS"),
        ("Annual value", f"R{s['annual_value']/1_000_000:,.1f}m", f"Target R{gauges['annual_benefits']['target']/1_000_000:,.1f}m"),
        ("Hours saved", f"{s['hours_saved']:,.0f}", f"Target {gauges['hours_saved']['target']:,.0f}"),
        ("Production use cases", str(s["production"]), f"Target {gauges['production_use_cases']['target']:,.0f}"),
        ("Cross-domain/market reusable", f"{s['reusable_percent']}%", f"{s['reusable_count']} of {s['active_use_cases']} use cases"),
        ("Portfolio ROI", f"{s['portfolio_roi']}%", f"Target {gauges['roi']['target']}%")
    ]
    for i, (label, value, note) in enumerate(kpis):
        add_kpi(slide, 0.45 + (i%3)*4.25, 1.25 + (i//3)*1.25, 3.75, .95, label, value, note)

    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "Executive Meter Gauges", f"Performance against {gauges['year']} targets")
    gauge_specs = [
        ("Annual benefits", gauges["annual_benefits"]["current"]/1_000_000, gauges["annual_benefits"]["target"]/1_000_000, "m", BASE_DIR / "static" / "gauge_value.png"),
        ("Hours returned", gauges["hours_saved"]["current"], gauges["hours_saved"]["target"], "", BASE_DIR / "static" / "gauge_hours.png"),
        ("Production adoption", gauges["production_use_cases"]["current"], gauges["production_use_cases"]["target"], "", BASE_DIR / "static" / "gauge_prod.png"),
        ("TMF maturity", gauges["tmf_maturity"]["current"], gauges["tmf_maturity"]["target"], "", BASE_DIR / "static" / "gauge_tmf.png"),
        ("Reusable use cases", gauges["reusable_percent"]["current"], gauges["reusable_percent"]["target"], "%", BASE_DIR / "static" / "gauge_reuse.png"),
        ("Portfolio ROI", gauges["roi"]["current"], gauges["roi"]["target"], "%", BASE_DIR / "static" / "gauge_roi.png")
    ]
    for title, current, target, unit, path in gauge_specs:
        make_gauge(current, target, title, unit, path)
    positions = [(0.35,1.1),(2.55,1.1),(4.75,1.1),(6.95,1.1),(9.15,1.1),(11.35,1.1)]
    for (_,_,_,_,path),(x,y) in zip(gauge_specs, positions):
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(1.9), Inches(1.55))

    tmp1 = BASE_DIR / "static" / "ppt_value_domain.png"
    tmp2 = BASE_DIR / "static" / "ppt_hours_domain.png"
    make_bar_chart(group_by(cases, "domain", "annual_value"), "Portfolio value by domain", "Rand", tmp1)
    make_bar_chart(group_by(cases, "domain", "hours_saved"), "Hours saved by domain", "Hours", tmp2)
    slide.shapes.add_picture(str(tmp1), Inches(.55), Inches(3.1), Inches(5.95), Inches(3.2))
    slide.shapes.add_picture(str(tmp2), Inches(6.85), Inches(3.1), Inches(5.95), Inches(3.2))

    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "AI Portfolio Heatmap", "Top use cases by annual value")

    # Build the heatmap table. Important:
    # PowerPoint table cells often default to a light/white fill.
    # If we only set the font to white, the data can appear invisible.
    # Therefore every cell gets an explicit dark fill and the header gets a red fill.
    top = sorted(cases, key=lambda x: as_number(x.get("annual_value", 0)), reverse=True)[:12]
    rows, cols = len(top) + 1, 8
    table_shape = slide.shapes.add_table(rows, cols, Inches(.25), Inches(1.05), Inches(12.85), Inches(5.95))
    table = table_shape.table

    headers = ["Use Case", "Domain", "TMF", "Impact", "Status", "KPI", "Value", "Reuse"]
    widths = [2.25, 1.35, .55, .75, .8, 4.25, .95, .5]
    for c, w in enumerate(widths):
        table.columns[c].width = Inches(w)

    def format_cell(cell, text, fill_rgb, font_rgb=(255, 255, 255), bold=False, size=7):
        cell.text = str(text or "")
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*fill_rgb)
        # Make margins tight so the table data stays visible on the slide.
        cell.margin_left = Inches(0.03)
        cell.margin_right = Inches(0.03)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(size)
            paragraph.font.bold = bold
            paragraph.font.color.rgb = RGBColor(*font_rgb)

    # Header row
    for c, h in enumerate(headers):
        format_cell(table.cell(0, c), h, (150, 0, 0), bold=True, size=7)

    # Body rows
    for r, item in enumerate(top, 1):
        status = item.get("status", "")
        impact = item.get("impact", "")
        reuse = str(item.get("reuse_score", ""))
        value = f"R{as_number(item.get('annual_value'))/1_000_000:.1f}m"
        vals = [
            item.get("name", ""),
            item.get("domain", ""),
            item.get("tmf_level", ""),
            impact,
            status,
            item.get("kpi", ""),
            value,
            reuse,
        ]
        for c, v in enumerate(vals):
            fill = (24, 28, 36) if r % 2 else (32, 37, 48)
            font = (255, 255, 255)
            # Subtle heatmap style by status / impact / reuse
            if c == 3 and "Very High" in impact:
                fill = (120, 20, 20)
            elif c == 3 and "High" in impact:
                fill = (110, 80, 10)
            elif c == 4 and status == "Production":
                fill = (20, 90, 45)
            elif c == 4 and status in ["Sandbox", "Prototype", "In Progress"]:
                fill = (25, 65, 95)
            elif c == 7 and as_number(reuse) >= 4:
                fill = (15, 85, 45)
            format_cell(table.cell(r, c), v, fill, font_rgb=font, bold=(c in [0, 6]), size=6.5)

    # Add a small footer note to clarify the table is generated from current dashboard data.
    note = slide.shapes.add_textbox(Inches(.28), Inches(7.05), Inches(12.5), Inches(.25))
    p = note.text_frame.paragraphs[0]
    p.text = "Source: Dashboard use case register. Values reflect current annual benefit estimates/configured data."
    p.font.size = Pt(7)
    p.font.color.rgb = RGBColor(180, 180, 185)

    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "Reuse Score Legend and Executive Decisions", "Build Once, Deploy Many governance")
    y = 1.15
    for r in REUSE_LEGEND:
        add_kpi(slide, .6, y, 5.2, .6, f"{r['score']} · {r['level']}", "", r["meaning"])
        y += .75
    y = 1.15
    for i, decision in enumerate(settings.get("executive_decisions", [])[:6], 1):
        add_kpi(slide, 6.3, y, 6.4, .6, f"Decision {i}", "", decision)
        y += .75

    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "Autonomous Network Journey", "TM Forum maturity progression mapped to ASOC capabilities")
    levels = [("L1","Assist","Ask ASOC, CoreXplore"),("L2","Augment","Ask Katlego, SOC Assistant"),("L3","Automate","Fault Localisation, Config Verification"),("L4","Closed Loop","Self-healing RAN, Agentic API"),("L5","Autonomous","Digital Twin, zero touch")]
    for i,(lvl,name,desc) in enumerate(levels):
        add_kpi(slide, .55+i*2.55, 2.05, 2.15, 1.55, lvl, name, desc)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return StreamingResponse(output, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Content-Disposition":"attachment; filename=asoc_ai_me_dashboard_export.pptx"})
