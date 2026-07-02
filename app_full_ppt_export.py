
from fastapi import FastAPI, Request, Form, UploadFile, File
from fastapi.responses import HTMLResponse, JSONResponse, RedirectResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from pathlib import Path
from collections import defaultdict
import json, csv, io, datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image, ImageDraw, ImageFont
import textwrap

BASE_DIR = Path(__file__).resolve().parent
DATA_FILE = BASE_DIR / "data" / "use_cases.json"
SETTINGS_FILE = BASE_DIR / "data" / "settings.json"

app = FastAPI(title="ASOC AI ME Dashboard", version="5.0.0")
app.mount("/static", StaticFiles(directory=str(BASE_DIR / "static")), name="static")
templates = Jinja2Templates(directory=str(BASE_DIR / "templates"))

FIELDS = [
    "id","name","domain","description","strategic_outcome","benefit_type","kpi",
    "baseline_kpi","current_kpi","target_kpi","benefit_confidence",
    "market","technical_owner","sponsor","status","tmf_level","impact","frequency",
    "annual_value","implementation_cost","hours_saved","mttr_improvement","reuse_score"
]

REUSE_LEGEND = [
    {"score": 1, "level": "Team Reuse", "meaning": "Works for one team or local process only."},
    {"score": 2, "level": "Function Reuse", "meaning": "Reusable within one ASOC function."},
    {"score": 3, "level": "Domain Reuse", "meaning": "Reusable across an entire FM, PM, CM or OSS domain."},
    {"score": 4, "level": "Cross-Domain Reuse", "meaning": "Reusable across multiple ASOC domains."},
    {"score": 5, "level": "Cross-Market Reuse", "meaning": "Build once, deploy across ZA, TZ, MZ, Lesotho and potentially Vodafone Group."},
]

def load_json(path, default):
    if not path.exists():
        return default
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

def save_json(path, data):
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)

def load_use_cases():
    return load_json(DATA_FILE, [])

def save_use_cases(data):
    save_json(DATA_FILE, data)

def load_settings():
    settings = load_json(SETTINGS_FILE, {})
    settings.setdefault("active_year", "FY27")
    settings.setdefault("targets", {})
    for year in ["FY27","FY28","FY29"]:
        settings["targets"].setdefault(year, {})
    return settings

def as_number(v, default=0):
    try:
        if v in [None, ""]:
            return default
        return float(v) if "." in str(v) else int(v)
    except Exception:
        return default

def safe_div(a, b):
    a = as_number(a); b = as_number(b)
    return 0 if b == 0 else a / b

def summary(use_cases, settings=None):
    settings = settings or load_settings()
    total_value = sum(as_number(x.get("annual_value", 0)) for x in use_cases)
    implementation_cost = sum(as_number(x.get("implementation_cost", 0)) for x in use_cases)
    total_hours = sum(as_number(x.get("hours_saved", 0)) for x in use_cases)
    production = sum(1 for x in use_cases if x.get("status") == "Production")
    sandbox = sum(1 for x in use_cases if x.get("status") == "Sandbox")
    prototype = sum(1 for x in use_cases if x.get("status") == "Prototype")
    in_progress = sum(1 for x in use_cases if x.get("status") == "In Progress")
    reusable_count = sum(1 for x in use_cases if as_number(x.get("reuse_score", 0)) >= 4)
    reusable_percent = round(safe_div(reusable_count, len(use_cases)) * 100, 1)
    roi = round(safe_div(total_value - implementation_cost, implementation_cost) * 100, 1) if implementation_cost else 0
    return {
        "active_use_cases": len(use_cases),
        "production": production,
        "sandbox": sandbox,
        "prototype": prototype,
        "in_progress": in_progress,
        "annual_value": total_value,
        "implementation_cost": implementation_cost,
        "hours_saved": total_hours,
        "avg_mttr": round(sum(as_number(x.get("mttr_improvement", 0)) for x in use_cases) / max(len(use_cases), 1), 1),
        "avg_reuse": round(sum(as_number(x.get("reuse_score", 0)) for x in use_cases) / max(len(use_cases), 1), 1),
        "reusable_count": reusable_count,
        "reusable_percent": reusable_percent,
        "portfolio_roi": roi,
        "tmf_current": as_number(settings.get("tmf_current", 1.8)),
        "tmf_target": as_number(settings.get("tmf_target", 3.0)),
    }

def active_targets(settings):
    active_year = settings.get("active_year", "FY27")
    targets = settings.get("targets", {}).get(active_year, {})
    return active_year, targets

def pct(current, target):
    return round(min(safe_div(current, target) * 100, 100), 1)

def gauge_values(s, settings):
    year, targets = active_targets(settings)
    return {
        "year": year,
        "annual_benefits": {
            "label": "Annual benefits realised",
            "current": s["annual_value"],
            "target": as_number(targets.get("annual_benefits", 1)),
            "unit": "R",
            "pct": pct(s["annual_value"], targets.get("annual_benefits", 1))
        },
        "hours_saved": {
            "label": "Hours returned",
            "current": s["hours_saved"],
            "target": as_number(targets.get("hours_saved", 1)),
            "unit": "hrs",
            "pct": pct(s["hours_saved"], targets.get("hours_saved", 1))
        },
        "production_use_cases": {
            "label": "Production adoption",
            "current": s["production"],
            "target": as_number(targets.get("production_use_cases", 1)),
            "unit": "use cases",
            "pct": pct(s["production"], targets.get("production_use_cases", 1))
        },
        "tmf_maturity": {
            "label": "TMF maturity",
            "current": s["tmf_current"],
            "target": as_number(targets.get("tmf_maturity", settings.get("tmf_target", 3.0))),
            "unit": "level",
            "pct": pct(s["tmf_current"], targets.get("tmf_maturity", settings.get("tmf_target", 3.0)))
        },
        "reusable_percent": {
            "label": "Cross-domain / market reusable",
            "current": s["reusable_percent"],
            "target": as_number(targets.get("reusable_use_cases_percent", 1)),
            "unit": "%",
            "pct": pct(s["reusable_percent"], targets.get("reusable_use_cases_percent", 1))
        },
        "roi": {
            "label": "Portfolio ROI",
            "current": s["portfolio_roi"],
            "target": as_number(targets.get("portfolio_roi_percent", 1)),
            "unit": "%",
            "pct": pct(s["portfolio_roi"], targets.get("portfolio_roi_percent", 1))
        }
    }

def group_by(use_cases, key, measure):
    buckets = defaultdict(float)
    for x in use_cases:
        buckets[x.get(key, "Unknown") or "Unknown"] += as_number(x.get(measure, 0))
    return [{"label": k, "value": v} for k, v in sorted(buckets.items(), key=lambda i: i[1], reverse=True)]

@app.get("/", response_class=HTMLResponse)
async def dashboard(request: Request):
    cases = load_use_cases()
    settings = load_settings()
    s = summary(cases, settings)
    return templates.TemplateResponse("index.html", {
        "request": request,
        "summary": s,
        "gauges": gauge_values(s, settings),
        "use_cases": cases,
        "settings": settings,
        "reuse_legend": REUSE_LEGEND
    })

@app.get("/admin", response_class=HTMLResponse)
async def admin(request: Request, edit: str = ""):
    cases = load_use_cases()
    settings = load_settings()
    record = next((x for x in cases if x.get("id") == edit), None)
    return templates.TemplateResponse("admin.html", {
        "request": request,
        "use_cases": cases,
        "settings": settings,
        "record": record,
        "fields": FIELDS,
        "reuse_legend": REUSE_LEGEND
    })

@app.post("/admin/save-use-case")
async def save_use_case(request: Request):
    form = await request.form()
    cases = load_use_cases()
    item = {}
    for field in FIELDS:
        value = form.get(field, "")
        if field in ["annual_value","implementation_cost","hours_saved","mttr_improvement","reuse_score"]:
            item[field] = as_number(value)
        else:
            item[field] = str(value).strip()
    if not item.get("id"):
        return RedirectResponse("/admin", status_code=303)
    found = False
    for i, c in enumerate(cases):
        if c.get("id") == item["id"]:
            cases[i] = item
            found = True
            break
    if not found:
        cases.append(item)
    save_use_cases(cases)
    return RedirectResponse("/admin", status_code=303)

@app.post("/admin/delete-use-case")
async def delete_use_case(id: str = Form(...)):
    cases = [x for x in load_use_cases() if x.get("id") != id]
    save_use_cases(cases)
    return RedirectResponse("/admin", status_code=303)

@app.post("/admin/save-settings")
async def save_settings(request: Request):
    form = await request.form()
    settings = load_settings()
    settings["dashboard_title"] = str(form.get("dashboard_title", settings.get("dashboard_title","ME Executive Dashboard")))
    settings["portfolio_period"] = str(form.get("portfolio_period", settings.get("portfolio_period","FY27–FY29")))
    settings["active_year"] = str(form.get("active_year", settings.get("active_year","FY27")))
    settings["tmf_current"] = as_number(form.get("tmf_current", settings.get("tmf_current", 1.8)))
    settings["tmf_target"] = as_number(form.get("tmf_target", settings.get("tmf_target", 3.0)))

    for year in ["FY27","FY28","FY29"]:
        settings.setdefault("targets", {}).setdefault(year, {})
        for metric in ["annual_benefits","hours_saved","production_use_cases","tmf_maturity","reusable_use_cases_percent","portfolio_roi_percent"]:
            settings["targets"][year][metric] = as_number(form.get(f"{year}_{metric}", settings["targets"][year].get(metric, 0)))

    decisions = str(form.get("executive_decisions", ""))
    settings["executive_decisions"] = [x.strip() for x in decisions.splitlines() if x.strip()]

    parsed_bets = []
    big_bets = str(form.get("big_bets", ""))
    for line in big_bets.splitlines():
        if "|" in line:
            name, desc = line.split("|", 1)
            parsed_bets.append({"name": name.strip(), "description": desc.strip()})
        elif line.strip():
            parsed_bets.append({"name": line.strip(), "description": ""})
    settings["big_bets"] = parsed_bets
    save_json(SETTINGS_FILE, settings)
    return RedirectResponse("/admin", status_code=303)

@app.get("/export/json")
async def export_json():
    data = {"use_cases": load_use_cases(), "settings": load_settings()}
    payload = json.dumps(data, indent=2).encode("utf-8")
    return StreamingResponse(io.BytesIO(payload), media_type="application/json", headers={"Content-Disposition":"attachment; filename=asoc_ai_dashboard_data.json"})

@app.get("/export/csv")
async def export_csv():
    output = io.StringIO()
    writer = csv.DictWriter(output, fieldnames=FIELDS)
    writer.writeheader()
    for row in load_use_cases():
        writer.writerow({k: row.get(k, "") for k in FIELDS})
    return StreamingResponse(io.BytesIO(output.getvalue().encode("utf-8")), media_type="text/csv", headers={"Content-Disposition":"attachment; filename=asoc_ai_use_cases.csv"})

@app.post("/import/csv")
async def import_csv(file: UploadFile = File(...)):
    content = (await file.read()).decode("utf-8-sig")
    reader = csv.DictReader(io.StringIO(content))
    existing_cases = load_use_cases()
    existing_by_id = {str(item.get("id", "")).strip(): item for item in existing_cases if str(item.get("id", "")).strip()}
    imported_by_id = {}
    for row in reader:
        item = {k: row.get(k, "") for k in FIELDS}
        item_id = str(item.get("id", "")).strip()
        if not item_id:
            continue
        item["id"] = item_id
        for n in ["annual_value","implementation_cost","hours_saved","mttr_improvement","reuse_score"]:
            item[n] = as_number(item.get(n, 0))
        imported_by_id[item_id] = item
    existing_by_id.update(imported_by_id)
    existing_order = [str(item.get("id", "")).strip() for item in existing_cases if str(item.get("id", "")).strip()]
    new_ids = [item_id for item_id in imported_by_id.keys() if item_id not in existing_order]
    final_ids = existing_order + new_ids
    merged_cases = [existing_by_id[item_id] for item_id in final_ids if item_id in existing_by_id]
    save_use_cases(merged_cases)
    return RedirectResponse("/admin", status_code=303)

@app.get("/api/use-cases")
async def api_use_cases():
    return JSONResponse(load_use_cases())

@app.get("/api/summary")
async def api_summary():
    cases = load_use_cases()
    settings = load_settings()
    s = summary(cases, settings)
    return JSONResponse({
        "summary": s,
        "gauges": gauge_values(s, settings),
        "value_by_domain": group_by(cases, "domain", "annual_value"),
        "hours_by_domain": group_by(cases, "domain", "hours_saved"),
        "value_by_sponsor": group_by(cases, "sponsor", "annual_value"),
        "status_mix": group_by(cases, "status", "annual_value"),
        "reuse_mix": group_by(cases, "reuse_score", "annual_value")
    })

def add_title(slide, title, subtitle=None):
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.5))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.42), Inches(0.78), Inches(12.2), Inches(0.3))
        p2 = st.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(220,220,220)

def set_bg(slide, color=(17,18,23)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)

def add_kpi(slide, x, y, w, h, label, value, note=""):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(34,37,48)
    shape.line.color.rgb = RGBColor(70,70,80)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(180,180,185)
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255,255,255)
    if note:
        p3 = tf.add_paragraph()
        p3.text = note
        p3.font.size = Pt(8)
        p3.font.color.rgb = RGBColor(170,170,175)

def make_bar_chart(data, title, ylabel, path):
    labels = [str(x["label"]) for x in data]
    vals = [x["value"] for x in data]
    fig, ax = plt.subplots(figsize=(7, 3.8), dpi=160)
    fig.patch.set_facecolor("#1a1c23")
    ax.set_facecolor("#1a1c23")
    ax.bar(labels, vals)
    ax.set_title(title, color="white", fontsize=12, pad=10)
    ax.set_ylabel(ylabel, color="white", fontsize=9)
    ax.tick_params(colors="white", labelsize=8)
    for spine in ax.spines.values():
        spine.set_color("#444")
    ax.grid(axis="y", color="white", alpha=.08)
    plt.xticks(rotation=20, ha="right")
    plt.tight_layout()
    fig.savefig(path, facecolor=fig.get_facecolor(), bbox_inches="tight")
    plt.close(fig)

def make_gauge(current, target, title, unit, path):
    current = float(current or 0)
    target = float(target or 1)
    pct_raw = (current / target) * 100 if target else 0
    pct_value = max(0, min(pct_raw / 100, 1))

    fig, ax = plt.subplots(figsize=(4.2, 2.8), dpi=180, subplot_kw={"projection": "polar"})
    fig.patch.set_facecolor("#1a1c23")
    ax.set_facecolor("#1a1c23")

    theta_bg = np.linspace(np.pi, 0, 100)
    ax.plot(theta_bg, [1]*len(theta_bg), linewidth=18, solid_capstyle="round", color="#e60000")

    theta_val = np.linspace(np.pi, np.pi - pct_value*np.pi, 100)
    ax.plot(theta_val, [1]*len(theta_val), linewidth=18, solid_capstyle="round", color="#22c55e")

    needle_theta = np.pi - pct_value*np.pi
    ax.plot([needle_theta, needle_theta], [0, .95], linewidth=2.5, color="white")
    ax.scatter([needle_theta], [0], s=30, color="white")

    ax.set_ylim(0, 1.25)
    ax.set_axis_off()

    current_label = f"{current:,.1f}" if current % 1 else f"{current:,.0f}"
    target_label = f"{target:,.1f}" if target % 1 else f"{target:,.0f}"
    pct_label = f"{pct_raw:,.0f}%"

    ax.text(.5, .42, pct_label, transform=ax.transAxes, ha="center", va="center",
            color="white", fontsize=18, fontweight="bold")
    ax.text(.5, .27, title, transform=ax.transAxes, ha="center", va="center",
            color="#d4d4d8", fontsize=9, fontweight="bold")
    ax.text(.5, .13, f"Current: {current_label}{unit}  |  Target: {target_label}{unit}",
            transform=ax.transAxes, ha="center", va="center", color="#a1a1aa", fontsize=7.5)

    plt.tight_layout(pad=.3)
    fig.savefig(path, facecolor=fig.get_facecolor(), bbox_inches="tight")
    plt.close(fig)



def _font(size=24, bold=False):
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ]
    for p in candidates:
        try:
            return ImageFont.truetype(p, size)
        except Exception:
            pass
    return ImageFont.load_default()

def _wrap(draw, text, font, width_px):
    text = str(text or "")
    words = text.split()
    lines, current = [], ""
    for w in words:
        test = (current + " " + w).strip()
        bbox = draw.textbbox((0, 0), test, font=font)
        if bbox[2] <= width_px or not current:
            current = test
        else:
            lines.append(current)
            current = w
    if current:
        lines.append(current)
    return lines

def _draw_wrapped(draw, xy, text, font, fill, width_px, max_lines=2, line_gap=2):
    x, y = xy
    lines = _wrap(draw, text, font, width_px)[:max_lines]
    for line in lines:
        draw.text((x, y), line, font=font, fill=fill)
        y += font.size + line_gap

def _color_for_status(status):
    s = str(status).lower()
    if 'production' in s or 'completed' in s:
        return (31, 181, 85)
    if 'progress' in s or 'prototype' in s:
        return (28, 130, 215)
    if 'planned' in s or 'not' in s:
        return (95, 105, 118)
    return (245, 177, 66)

def _color_for_impact(impact):
    s = str(impact).lower()
    if 'very' in s:
        return (230, 0, 0)
    if 'high' in s:
        return (255, 176, 0)
    if 'medium' in s:
        return (255, 220, 0)
    return (104, 196, 67)

def _color_for_reuse(reuse):
    try:
        r = int(float(reuse))
    except Exception:
        r = 0
    if r >= 5:
        return (24, 200, 74)
    if r == 4:
        return (77, 195, 105)
    if r == 3:
        return (255, 214, 0)
    if r == 2:
        return (255, 176, 0)
    return (240, 52, 40)

def make_heatmap_slide_image(cases, path):
    """Create a self-contained heatmap image so the PPT slide always renders visibly."""
    W, H = 1920, 1080
    img = Image.new('RGB', (W, H), (6, 19, 28))
    draw = ImageDraw.Draw(img)
    # subtle background
    for i in range(0, H, 5):
        shade = int(18 + 18 * (i / H))
        draw.line([(0, i), (W, i)], fill=(4, shade, 30))
    red = (230, 0, 0)
    white = (245, 247, 250)
    muted = (175, 186, 196)
    panel = (10, 30, 42)
    line = (72, 105, 125)

    title_f = _font(48, True)
    sub_f = _font(22, False)
    hdr_f = _font(20, True)
    cell_f = _font(17, False)
    cell_b = _font(17, True)
    small_f = _font(15, False)

    draw.text((60, 42), "AI Portfolio Heatmap", font=title_f, fill=white)
    draw.text((62, 102), "Top use cases by annual value · generated from dashboard data", font=sub_f, fill=muted)
    # Vodacom dot
    draw.ellipse((1780, 42, 1850, 112), outline=red, width=12)
    draw.text((1695, 119), "ASOC AI", font=hdr_f, fill=white)

    top = sorted(cases, key=lambda x: as_number(x.get("annual_value",0)), reverse=True)[:12]
    x0, y0 = 45, 170
    table_w = 1830
    header_h = 48
    row_h = 64
    # Columns: use case, domain, tmf, impact, status, kpi, value, reuse
    widths = [360, 230, 95, 130, 135, 540, 150, 95]
    headers = ["Use Case", "Domain", "TMF", "Impact", "Status", "KPI / Benefit Driver", "Value", "Reuse"]
    xs = [x0]
    for w in widths[:-1]:
        xs.append(xs[-1] + w)

    # table panel
    draw.rounded_rectangle((x0-10, y0-12, x0+table_w+10, y0+header_h+row_h*len(top)+20), radius=18, fill=(8, 24, 36), outline=line, width=2)
    draw.rectangle((x0, y0, x0+table_w, y0+header_h), fill=red)
    for i, h in enumerate(headers):
        draw.text((xs[i]+10, y0+13), h, font=hdr_f, fill=white)

    for r, item in enumerate(top):
        y = y0 + header_h + r * row_h
        fill = (13, 34, 48) if r % 2 == 0 else (9, 28, 40)
        draw.rectangle((x0, y, x0+table_w, y+row_h), fill=fill)
        draw.line((x0, y, x0+table_w, y), fill=(48, 74, 92), width=1)
        vals = [
            item.get("name", ""), item.get("domain", ""), item.get("tmf_level", ""), item.get("impact", ""), item.get("status", ""),
            item.get("kpi", ""), f"R{as_number(item.get('annual_value'))/1_000_000:.1f}m", str(item.get("reuse_score", ""))
        ]
        for c, val in enumerate(vals):
            x = xs[c]
            w = widths[c]
            if c == 3:  # impact pill
                col = _color_for_impact(val)
                draw.rounded_rectangle((x+8, y+15, x+w-8, y+47), radius=12, fill=col)
                draw.text((x+16, y+22), str(val)[:12], font=small_f, fill=(0,0,0) if col[0] > 220 and col[1] > 170 else white)
            elif c == 4:  # status pill
                col = _color_for_status(val)
                draw.rounded_rectangle((x+8, y+15, x+w-8, y+47), radius=12, fill=col)
                _draw_wrapped(draw, (x+14, y+21), val, small_f, white, w-25, max_lines=1)
            elif c == 7:  # reuse score
                col = _color_for_reuse(val)
                draw.ellipse((x+28, y+13, x+64, y+49), fill=col)
                draw.text((x+40, y+20), str(val), font=cell_b, fill=(0,0,0))
            elif c == 6:
                draw.text((x+10, y+22), val, font=cell_b, fill=(31, 255, 83))
            else:
                _draw_wrapped(draw, (x+10, y+11), val, cell_f if c != 0 else cell_b, white if c != 1 else muted, w-18, max_lines=2)
        # vertical lines
        xx = x0
        for w in widths:
            draw.line((xx, y, xx, y+row_h), fill=(42, 65, 82), width=1)
            xx += w

    # footer legend
    fy = y0 + header_h + row_h*len(top) + 45
    draw.text((60, fy), "Legend:", font=hdr_f, fill=white)
    lx = 165
    for label, col in [("Very High / High Impact", (255,176,0)), ("Production/Completed", (31,181,85)), ("Sandbox/In Progress", (28,130,215)), ("Cross-Market Reuse", (24,200,74))]:
        draw.rounded_rectangle((lx, fy-2, lx+22, fy+20), radius=5, fill=col)
        draw.text((lx+32, fy-4), label, font=small_f, fill=muted)
        lx += 330
    img.save(path)



def _fmt_money(v):
    try:
        return f"R{float(v)/1_000_000:.1f}m"
    except Exception:
        return "R0.0m"

def make_table_slide_image(title, subtitle, headers, rows, col_widths, path, max_lines=None):
    """Create a full-slide table image. This avoids PowerPoint native table rendering issues and supports many slides."""
    W, H = 1920, 1080
    img = Image.new('RGB', (W, H), (6, 19, 28))
    draw = ImageDraw.Draw(img)
    for i in range(0, H, 6):
        shade = int(16 + 18 * (i / H))
        draw.line([(0, i), (W, i)], fill=(4, shade, 30))

    red = (230, 0, 0)
    white = (245, 247, 250)
    muted = (175, 186, 196)
    line = (72, 105, 125)

    title_f = _font(46, True)
    sub_f = _font(21, False)
    hdr_f = _font(20, True)
    cell_f = _font(17, False)
    cell_b = _font(17, True)
    small_f = _font(15, False)

    draw.text((55, 38), title, font=title_f, fill=white)
    if subtitle:
        draw.text((57, 96), subtitle, font=sub_f, fill=muted)
    draw.ellipse((1780, 42, 1850, 112), outline=red, width=12)
    draw.text((1695, 119), "ASOC AI", font=hdr_f, fill=white)

    x0, y0 = 45, 158
    table_w = sum(col_widths)
    header_h = 48
    available_h = H - y0 - 80
    row_h = int((available_h - header_h) / max(len(rows), 1)) if rows else 64
    row_h = max(48, min(82, row_h))
    max_rows_fit = int((available_h - header_h) / row_h)
    rows = rows[:max_rows_fit]

    xs = [x0]
    for w in col_widths[:-1]:
        xs.append(xs[-1] + w)

    draw.rounded_rectangle((x0-10, y0-12, x0+table_w+10, y0+header_h+row_h*len(rows)+20), radius=18, fill=(8, 24, 36), outline=line, width=2)
    draw.rectangle((x0, y0, x0+table_w, y0+header_h), fill=red)
    for i, h in enumerate(headers):
        _draw_wrapped(draw, (xs[i]+10, y0+12), h, hdr_f, white, col_widths[i]-18, max_lines=1)

    for r, row in enumerate(rows):
        y = y0 + header_h + r*row_h
        fill = (13, 34, 48) if r % 2 == 0 else (9, 28, 40)
        draw.rectangle((x0, y, x0+table_w, y+row_h), fill=fill)
        draw.line((x0, y, x0+table_w, y), fill=(48, 74, 92), width=1)
        xx = x0
        for w in col_widths:
            draw.line((xx, y, xx, y+row_h), fill=(42,65,82), width=1)
            xx += w
        for c, val in enumerate(row):
            x = xs[c]; w = col_widths[c]
            val = '' if val is None else str(val)
            # special formatting on status, impact and reuse where header names match
            header = headers[c].lower()
            if 'status' in header:
                col = _color_for_status(val)
                draw.rounded_rectangle((x+8, y+12, x+w-8, y+42), radius=12, fill=col)
                _draw_wrapped(draw, (x+14, y+18), val, small_f, white, w-25, max_lines=1)
            elif 'impact' in header:
                col = _color_for_impact(val)
                draw.rounded_rectangle((x+8, y+12, x+w-8, y+42), radius=12, fill=col)
                _draw_wrapped(draw, (x+14, y+18), val, small_f, (0,0,0) if col[1] > 170 else white, w-25, max_lines=1)
            elif 'reuse' in header:
                col = _color_for_reuse(val)
                draw.ellipse((x+20, y+10, x+54, y+44), fill=col)
                draw.text((x+31, y+17), val, font=cell_b, fill=(0,0,0))
            elif 'value' in header or 'benefit' in header or 'roi' in header:
                _draw_wrapped(draw, (x+10, y+14), val, cell_b, (31, 255, 83), w-18, max_lines=max_lines or 2)
            else:
                _draw_wrapped(draw, (x+10, y+10), val, cell_b if c == 0 else cell_f, white if c != 1 else muted, w-18, max_lines=max_lines or 2)
    # footer
    draw.text((55, H-48), "Source: ASOC AI dashboard use case register · Exported from live dashboard", font=small_f, fill=muted)
    img.save(path)

def make_targets_slide_image(settings, summary_data, gauges, path):
    rows = []
    targets = settings.get('targets', {})
    metrics = [
        ('annual_benefits', 'Annual Benefits', lambda v: _fmt_money(v)),
        ('hours_saved', 'Hours Saved', lambda v: f"{as_number(v):,.0f}"),
        ('production_use_cases', 'Production Use Cases', lambda v: f"{as_number(v):,.0f}"),
        ('tmf_maturity', 'TMF Maturity', lambda v: f"L{as_number(v):.1f}"),
        ('reusable_use_cases_percent', 'Reusable Use Cases %', lambda v: f"{as_number(v):.0f}%"),
        ('portfolio_roi_percent', 'Portfolio ROI %', lambda v: f"{as_number(v):.0f}%"),
    ]
    for key, label, fmt in metrics:
        rows.append([label, fmt(targets.get('FY27', {}).get(key, 0)), fmt(targets.get('FY28', {}).get(key, 0)), fmt(targets.get('FY29', {}).get(key, 0))])
    make_table_slide_image(
        "Targets and Gauge Configuration",
        f"Active year: {settings.get('active_year','FY27')} · Portfolio value: {_fmt_money(summary_data.get('annual_value',0))}",
        ["Metric", "FY27", "FY28", "FY29"],
        rows,
        [620, 380, 380, 380],
        path,
        max_lines=1,
    )

def chunked(items, size):
    for i in range(0, len(items), size):
        yield i, items[i:i+size]

@app.get("/export/pptx")
async def export_pptx():
    cases = load_use_cases()
    settings = load_settings()
    s = summary(cases, settings)
    gauges = gauge_values(s, settings)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Executive summary
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "ASOC AI ME Executive Dashboard", f"{settings.get('active_year','FY27')} · Exported {datetime.date.today().isoformat()}")
    kpis = [
        ("Active AI use cases", str(s["active_use_cases"]), "FM · PM · CM · OSS"),
        ("Annual value", f"R{s['annual_value']/1_000_000:,.1f}m", f"Target R{gauges['annual_benefits']['target']/1_000_000:,.1f}m"),
        ("Hours saved", f"{s['hours_saved']:,.0f}", f"Target {gauges['hours_saved']['target']:,.0f}"),
        ("Production use cases", str(s["production"]), f"Target {gauges['production_use_cases']['target']:,.0f}"),
        ("Cross-domain/market reusable", f"{s['reusable_percent']}%", f"{s['reusable_count']} of {s['active_use_cases']} use cases"),
        ("Portfolio ROI", f"{s['portfolio_roi']}%", f"Target {gauges['roi']['target']}%")
    ]
    for i, (label, value, note) in enumerate(kpis):
        add_kpi(slide, 0.45 + (i%3)*4.25, 1.25 + (i//3)*1.25, 3.75, .95, label, value, note)

    # Slide 2: Gauges + charts
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "Executive Meter Gauges", f"Performance against {gauges['year']} targets")
    gauge_specs = [
        ("Annual benefits", gauges["annual_benefits"]["current"]/1_000_000, gauges["annual_benefits"]["target"]/1_000_000, "m", BASE_DIR / "static" / "gauge_value.png"),
        ("Hours returned", gauges["hours_saved"]["current"], gauges["hours_saved"]["target"], "", BASE_DIR / "static" / "gauge_hours.png"),
        ("Production adoption", gauges["production_use_cases"]["current"], gauges["production_use_cases"]["target"], "", BASE_DIR / "static" / "gauge_prod.png"),
        ("TMF maturity", gauges["tmf_maturity"]["current"], gauges["tmf_maturity"]["target"], "", BASE_DIR / "static" / "gauge_tmf.png"),
        ("Reusable use cases", gauges["reusable_percent"]["current"], gauges["reusable_percent"]["target"], "%", BASE_DIR / "static" / "gauge_reuse.png"),
        ("Portfolio ROI", gauges["roi"]["current"], gauges["roi"]["target"], "%", BASE_DIR / "static" / "gauge_roi.png")
    ]
    for title, current, target, unit, path in gauge_specs:
        make_gauge(current, target, title, unit, path)
    positions = [(0.35,1.1),(2.55,1.1),(4.75,1.1),(6.95,1.1),(9.15,1.1),(11.35,1.1)]
    for (_,_,_,_,path),(x,y) in zip(gauge_specs, positions):
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(1.9), Inches(1.55))
    tmp1 = BASE_DIR / "static" / "ppt_value_domain.png"
    tmp2 = BASE_DIR / "static" / "ppt_hours_domain.png"
    make_bar_chart(group_by(cases, "domain", "annual_value"), "Portfolio value by domain", "Rand", tmp1)
    make_bar_chart(group_by(cases, "domain", "hours_saved"), "Hours saved by domain", "Hours", tmp2)
    slide.shapes.add_picture(str(tmp1), Inches(.55), Inches(3.1), Inches(5.95), Inches(3.2))
    slide.shapes.add_picture(str(tmp2), Inches(6.85), Inches(3.1), Inches(5.95), Inches(3.2))

    # Slide 3: Top heatmap
    heatmap_png = BASE_DIR / "static" / "ppt_heatmap.png"
    make_heatmap_slide_image(cases, heatmap_png)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(heatmap_png), Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    # Slide 4: Targets / gauge configuration
    targets_png = BASE_DIR / "static" / "ppt_targets.png"
    make_targets_slide_image(settings, s, gauges, targets_png)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(targets_png), Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    # Slides 5+: Complete use case register, split over multiple slides
    sorted_cases = sorted(cases, key=lambda x: as_number(x.get("annual_value",0)), reverse=True)
    for idx, chunk in chunked(sorted_cases, 8):
        rows = []
        for item in chunk:
            rows.append([
                item.get("id", ""),
                item.get("name", ""),
                item.get("domain", ""),
                item.get("status", ""),
                item.get("kpi", ""),
                _fmt_money(as_number(item.get("annual_value", 0))),
                str(item.get("reuse_score", "")),
            ])
        out = BASE_DIR / "static" / f"ppt_register_{idx}.png"
        make_table_slide_image(
            "Complete AI Use Case Register",
            f"All dashboard records · rows {idx+1}-{idx+len(chunk)} of {len(sorted_cases)}",
            ["ID", "Use Case", "Domain", "Status", "KPI / Benefit Driver", "Value", "Reuse"],
            rows,
            [140, 390, 260, 160, 600, 150, 100],
            out,
            max_lines=3,
        )
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(out), Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    # Domain detail slides: FM, PM, CM, OSS
    domains = []
    for c in cases:
        if c.get("domain") not in domains:
            domains.append(c.get("domain"))
    for domain in domains:
        domain_cases = [c for c in sorted_cases if c.get("domain") == domain]
        rows = []
        for item in domain_cases:
            rows.append([
                item.get("name", ""), item.get("tmf_level", ""), item.get("impact", ""), item.get("status", ""),
                item.get("sponsor", ""), item.get("technical_owner", ""), _fmt_money(as_number(item.get("annual_value",0))), str(item.get("reuse_score", ""))
            ])
        out = BASE_DIR / "static" / f"ppt_domain_{domain.replace(' ','_')}.png"
        make_table_slide_image(
            f"Domain Detail · {domain}",
            f"{len(domain_cases)} use cases · total value {_fmt_money(sum(as_number(x.get('annual_value',0)) for x in domain_cases))}",
            ["Use Case", "TMF", "Impact", "Status", "Sponsor", "Tech Owner", "Value", "Reuse"],
            rows,
            [390, 95, 130, 150, 220, 260, 140, 90],
            out,
            max_lines=2,
        )
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(out), Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    # KPI detail slides: baseline/current/target/strategic outcome
    for idx, chunk in chunked(sorted_cases, 6):
        rows = []
        for item in chunk:
            rows.append([
                item.get("id", ""), item.get("name", ""), item.get("strategic_outcome", ""),
                item.get("baseline_kpi", ""), item.get("current_kpi", ""), item.get("target_kpi", ""), item.get("benefit_confidence", "")
            ])
        out = BASE_DIR / "static" / f"ppt_kpi_detail_{idx}.png"
        make_table_slide_image(
            "KPI Detail and Strategic Outcomes",
            f"Baseline, current and target KPI tracking · rows {idx+1}-{idx+len(chunk)} of {len(sorted_cases)}",
            ["ID", "Use Case", "Strategic Outcome", "Baseline KPI", "Current KPI", "Target KPI", "Confidence"],
            rows,
            [110, 300, 420, 270, 270, 330, 120],
            out,
            max_lines=3,
        )
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(out), Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    # Reuse score and decisions
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "Reuse Score Legend and Executive Decisions", "Build Once, Deploy Many governance")
    y = 1.15
    for r in REUSE_LEGEND:
        add_kpi(slide, .6, y, 5.2, .6, f"{r['score']} · {r['level']}", "", r["meaning"])
        y += .75
    y = 1.15
    for i, decision in enumerate(settings.get("executive_decisions", [])[:6], 1):
        add_kpi(slide, 6.3, y, 6.4, .6, f"Decision {i}", "", decision)
        y += .75

    # Autonomous journey
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "Autonomous Network Journey", "TM Forum maturity progression mapped to ASOC capabilities")
    levels = [("L1","Assist","Ask ASOC, CoreXplore"),("L2","Augment","Ask Katlego, SOC Assistant"),("L3","Automate","Fault Localisation, Config Verification"),("L4","Closed Loop","Self-healing RAN, Agentic API"),("L5","Autonomous","Digital Twin, zero touch")]
    for i,(lvl,name,desc) in enumerate(levels):
        add_kpi(slide, .55+i*2.55, 2.05, 2.15, 1.55, lvl, name, desc)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return StreamingResponse(output, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Content-Disposition":"attachment; filename=asoc_ai_me_dashboard_export.pptx"})
